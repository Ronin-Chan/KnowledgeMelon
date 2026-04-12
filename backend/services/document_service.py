import csv
import hashlib
import json
import os
import asyncio
import shutil
import subprocess
import tempfile
import xml.etree.ElementTree as ET
from typing import List
from pypdf import PdfReader
from docx import Document as DocxDocument
import uuid

from bs4 import BeautifulSoup

from services.file_types import SUPPORTED_FILE_EXTENSION_SET

try:
    import fitz
except ImportError:
    fitz = None

try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    import pymupdf4llm
except ImportError:
    pymupdf4llm = None

try:
    from rapidocr_onnxruntime import RapidOCR
except ImportError:
    RapidOCR = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import yaml
except ImportError:
    yaml = None

try:
    from striprtf.striprtf import rtf_to_text
except ImportError:
    rtf_to_text = None

try:
    import xlrd
except ImportError:
    xlrd = None

MEANINGFUL_TEXT_MIN_CHARS = 60
MAX_TABLE_ROWS = 1000
MAX_SHEET_ROWS = 1000


class DocumentService:
    def __init__(self, upload_dir: str = "uploads"):
        self.upload_dir = upload_dir
        os.makedirs(upload_dir, exist_ok=True)

    async def save_document(
        self,
        file_name: str,
        content: bytes,
        owner_id: str | None = None,
    ) -> str:
        """Save uploaded document and return file path"""
        file_id = str(uuid.uuid4())
        file_ext = os.path.splitext(file_name)[1].lower()
        target_dir = os.path.join(self.upload_dir, owner_id) if owner_id else self.upload_dir
        os.makedirs(target_dir, exist_ok=True)
        file_path = os.path.join(target_dir, f"{file_id}{file_ext}")

        with open(file_path, "wb") as f:
            f.write(content)

        return file_path, file_ext

    def compute_content_hash(self, content: bytes) -> str:
        hasher = hashlib.sha256()
        hasher.update(content)
        return hasher.hexdigest()

    def build_document_metadata(
        self,
        file_name: str,
        file_type: str,
        content_hash: str,
        extra: dict | None = None,
    ) -> str:
        payload = {
            "source_name": file_name,
            "file_type": file_type,
            "content_hash": content_hash,
        }
        if extra:
            payload.update(extra)
        return json.dumps(payload, ensure_ascii=False)

    def chunk_hash(self, chunk_text: str) -> str:
        hasher = hashlib.sha256()
        hasher.update(chunk_text.encode("utf-8"))
        return hasher.hexdigest()

    async def parse_document(
        self,
        file_path: str,
        file_type: str,
        start_page: int = None,
        end_page: int = None,
    ) -> str:
        """Parse document and return text content

        Args:
            file_path: Path to the document
            file_type: File extension
            start_page: For PDFs, start from this page (0-indexed, optional)
            end_page: For PDFs, end at this page (exclusive, optional)
        """
        file_type = (file_type or "").lower()

        if file_type not in SUPPORTED_FILE_EXTENSION_SET:
            raise ValueError(f"Unsupported file type: {file_type}")

        if file_type == ".pdf":
            return await asyncio.to_thread(
                self._parse_pdf,
                file_path,
                start_page or 0,
                end_page,
            )
        elif file_type == ".docx":
            return await asyncio.to_thread(self._parse_docx, file_path)
        elif file_type in [".doc", ".ppt"]:
            return await asyncio.to_thread(self._parse_binary_office, file_path, file_type)
        elif file_type == ".pptx":
            return await asyncio.to_thread(self._parse_pptx, file_path)
        elif file_type == ".xlsx":
            return await asyncio.to_thread(self._parse_xlsx, file_path)
        elif file_type == ".xls":
            return await asyncio.to_thread(self._parse_xls, file_path)
        elif file_type == ".csv":
            return await asyncio.to_thread(self._parse_csv, file_path, ",")
        elif file_type == ".tsv":
            return await asyncio.to_thread(self._parse_csv, file_path, "\t")
        elif file_type == ".json":
            return await asyncio.to_thread(self._parse_json, file_path)
        elif file_type == ".xml":
            return await asyncio.to_thread(self._parse_xml, file_path)
        elif file_type in [".yaml", ".yml"]:
            return await asyncio.to_thread(self._parse_yaml, file_path)
        elif file_type in [".html", ".htm"]:
            return await asyncio.to_thread(self._parse_html, file_path)
        elif file_type == ".rtf":
            return await asyncio.to_thread(self._parse_rtf, file_path)
        elif file_type in [".log", ".srt", ".vtt"]:
            return await asyncio.to_thread(self._parse_text, file_path)
        elif file_type in [".png", ".jpg", ".jpeg", ".webp"]:
            return await asyncio.to_thread(self._parse_image, file_path)
        elif file_type in [".txt", ".md", ".py", ".js", ".ts", ".java", ".go"]:
            return await asyncio.to_thread(self._parse_text, file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

    def _parse_pdf(
        self, file_path: str, start_page: int = 0, end_page: int = None
    ) -> str:
        extractor_sequence = [
            self._parse_pdf_with_pymupdf4llm,
            self._parse_pdf_with_pypdf,
            self._parse_pdf_with_pymupdf,
            self._parse_pdf_with_pdfplumber,
            self._parse_pdf_with_rapidocr,
        ]

        first_non_empty_text = ""

        for extractor in extractor_sequence:
            try:
                extracted_text = extractor(file_path, start_page, end_page)
            except Exception:
                continue

            if extracted_text and extracted_text.strip() and not first_non_empty_text:
                first_non_empty_text = extracted_text

            if self.has_meaningful_text(extracted_text):
                return extracted_text

        return first_non_empty_text

    def _parse_pdf_with_pymupdf4llm(
        self, file_path: str, start_page: int = 0, end_page: int = None
    ) -> str:
        if pymupdf4llm is None:
            return ""

        doc = fitz.open(file_path) if fitz is not None else None
        total_pages = doc.page_count if doc is not None else 0
        if doc is not None:
            doc.close()

        if total_pages <= 0:
            return ""

        start = max(0, start_page)
        end = min(total_pages, end_page) if end_page is not None else total_pages
        if end <= start:
            return ""

        pages = list(range(start, end))
        markdown_text = pymupdf4llm.to_markdown(file_path, pages=pages)
        return markdown_text if isinstance(markdown_text, str) else ""

    def _parse_pdf_pages_with_pymupdf4llm(
        self, file_path: str, start_page: int = 0, end_page: int = None
    ) -> list[str]:
        if pymupdf4llm is None:
            return []

        doc = fitz.open(file_path) if fitz is not None else None
        total_pages = doc.page_count if doc is not None else 0
        if doc is not None:
            doc.close()

        if total_pages <= 0:
            return []

        start = max(0, start_page)
        end = min(total_pages, end_page) if end_page is not None else total_pages
        if end <= start:
            return []

        pages: list[str] = []
        for page_index in range(start, end):
            page_markdown = pymupdf4llm.to_markdown(file_path, pages=[page_index])
            if isinstance(page_markdown, str):
                pages.append(page_markdown)
        return pages

    def _parse_pdf_with_pypdf(
        self, file_path: str, start_page: int = 0, end_page: int = None
    ) -> str:
        reader = PdfReader(file_path)
        total_pages = len(reader.pages)
        start = max(0, start_page)
        end = min(total_pages, end_page) if end_page is not None else total_pages

        text_parts: list[str] = []
        for i in range(start, end):
            page = reader.pages[i]
            page_text = page.extract_text()
            if page_text:
                text_parts.append(f"\n--- 第 {i + 1} 页 ---\n{page_text}\n")

        return "".join(text_parts)

    def _parse_pdf_with_pymupdf(
        self, file_path: str, start_page: int = 0, end_page: int = None
    ) -> str:
        if fitz is None:
            return ""

        doc = fitz.open(file_path)
        try:
            total_pages = doc.page_count
            start = max(0, start_page)
            end = min(total_pages, end_page) if end_page is not None else total_pages

            text_parts: list[str] = []
            for i in range(start, end):
                page = doc.load_page(i)
                page_text = page.get_text("text")
                if page_text:
                    text_parts.append(f"\n--- 第 {i + 1} 页 ---\n{page_text}\n")

            return "".join(text_parts)
        finally:
            doc.close()

    def _parse_pdf_with_pdfplumber(
        self, file_path: str, start_page: int = 0, end_page: int = None
    ) -> str:
        if pdfplumber is None:
            return ""

        with pdfplumber.open(file_path) as pdf:
            total_pages = len(pdf.pages)
            start = max(0, start_page)
            end = min(total_pages, end_page) if end_page is not None else total_pages

            text_parts: list[str] = []
            for i in range(start, end):
                page = pdf.pages[i]
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(f"\n--- 第 {i + 1} 页 ---\n{page_text}\n")

            return "".join(text_parts)

    def _parse_pdf_with_rapidocr(
        self, file_path: str, start_page: int = 0, end_page: int = None
    ) -> str:
        if RapidOCR is None or fitz is None:
            return ""

        ocr_engine = RapidOCR()
        doc = fitz.open(file_path)
        try:
            total_pages = doc.page_count
            start = max(0, start_page)
            end = min(total_pages, end_page) if end_page is not None else total_pages

            text_parts: list[str] = []
            for i in range(start, end):
                page = doc.load_page(i)
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)

                temp_path = ""
                try:
                    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
                        temp_path = f.name
                        pix.save(temp_path)

                    ocr_result, _ = ocr_engine(temp_path)
                    if not ocr_result:
                        continue

                    line_texts: list[str] = []
                    for item in ocr_result:
                        if isinstance(item, list) and len(item) >= 2:
                            line = item[1]
                            if isinstance(line, str) and line.strip():
                                line_texts.append(line.strip())

                    if line_texts:
                        joined_text = "\n".join(line_texts)
                        text_parts.append(
                            f"\n--- 第 {i + 1} 页 ---\n{joined_text}\n"
                        )
                finally:
                    if temp_path and os.path.exists(temp_path):
                        os.remove(temp_path)

            return "".join(text_parts)
        finally:
            doc.close()

    def has_meaningful_text(
        self, text: str, min_chars: int = MEANINGFUL_TEXT_MIN_CHARS
    ) -> bool:
        compact = "".join(text.split())
        valid_chars = sum(1 for ch in compact if ch.isalnum())
        return valid_chars >= min_chars

    def get_pdf_info(self, file_path: str) -> dict:
        """Get PDF metadata (page count, file size)"""
        reader = PdfReader(file_path)
        file_size = os.path.getsize(file_path)
        return {
            "total_pages": len(reader.pages),
            "file_size_bytes": file_size,
            "file_size_mb": round(file_size / (1024 * 1024), 2),
        }

    def _parse_docx(self, file_path: str) -> str:
        doc = DocxDocument(file_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text

    def _parse_binary_office(self, file_path: str, file_type: str) -> str:
        command_name = "catppt" if file_type == ".ppt" else "catdoc"
        executable = shutil.which(command_name)
        if executable is None:
            raise ValueError(
                f"{file_type.upper()} parsing requires {command_name} to be installed"
            )

        try:
            result = subprocess.run(
                [executable, file_path],
                capture_output=True,
                text=True,
                check=False,
            )
        except Exception as exc:
            raise ValueError(f"Failed to parse {file_type.upper()} file: {exc}") from exc

        output = (result.stdout or result.stderr or "").strip()
        if not output:
            raise ValueError(f"No text extracted from {file_type.upper()} file")
        return output

    def _parse_pptx(self, file_path: str) -> str:
        if Presentation is None:
            raise ValueError("PPTX parsing requires python-pptx to be installed")

        prs = Presentation(file_path)
        parts: list[str] = []
        for slide_index, slide in enumerate(prs.slides, start=1):
            slide_parts: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = shape.text_frame.text.strip()
                    if text:
                        slide_parts.append(text)
                if getattr(shape, "has_table", False):
                    table_rows: list[list[str]] = []
                    for row in shape.table.rows:
                        table_rows.append([self._normalize_cell_value(cell.text) for cell in row.cells])
                    if table_rows:
                        slide_parts.append(self._format_table_rows(table_rows))
            if slide_parts:
                parts.append(f"--- Slide {slide_index} ---\n" + "\n".join(slide_parts))
        return "\n\n".join(parts).strip()

    def _parse_xlsx(self, file_path: str) -> str:
        if load_workbook is None:
            raise ValueError("XLSX parsing requires openpyxl to be installed")

        workbook = load_workbook(file_path, read_only=True, data_only=True)
        parts: list[str] = []
        try:
            for sheet in workbook.worksheets:
                sheet_rows = []
                for row_index, row in enumerate(sheet.iter_rows(values_only=True)):
                    if row_index >= MAX_SHEET_ROWS:
                        break
                    sheet_rows.append([self._normalize_cell_value(cell) for cell in row])
                if sheet_rows:
                    parts.append(f"--- Sheet: {sheet.title} ---\n{self._format_table_rows(sheet_rows)}")
        finally:
            workbook.close()
        return "\n\n".join(parts).strip()

    def _parse_xls(self, file_path: str) -> str:
        if xlrd is None:
            raise ValueError("XLS parsing requires xlrd to be installed")

        workbook = xlrd.open_workbook(file_path)
        parts: list[str] = []
        for sheet in workbook.sheets():
            sheet_rows = []
            for row_index in range(min(sheet.nrows, MAX_SHEET_ROWS)):
                row = [
                    self._normalize_cell_value(sheet.cell_value(row_index, col_index))
                    for col_index in range(sheet.ncols)
                ]
                sheet_rows.append(row)
            if sheet_rows:
                parts.append(f"--- Sheet: {sheet.name} ---\n{self._format_table_rows(sheet_rows)}")
        return "\n\n".join(parts).strip()

    def _parse_csv(self, file_path: str, delimiter: str) -> str:
        parts: list[str] = []
        with open(file_path, "r", encoding="utf-8-sig", errors="ignore", newline="") as f:
            reader = csv.reader(f, delimiter=delimiter)
            for row_index, row in enumerate(reader):
                if row_index >= MAX_TABLE_ROWS:
                    break
                parts.append("\t".join(cell.strip() for cell in row))
        return "\n".join(parts).strip()

    def _parse_json(self, file_path: str) -> str:
        raw = self._parse_text(file_path)
        try:
            data = json.loads(raw)
        except Exception:
            return raw
        return json.dumps(data, ensure_ascii=False, indent=2)

    def _parse_xml(self, file_path: str) -> str:
        raw = self._parse_text(file_path)
        try:
            root = ET.fromstring(raw)
        except Exception:
            return raw

        parts: list[str] = []
        for elem in root.iter():
            text = " ".join((elem.text or "").split())
            if text:
                tag_name = elem.tag.split("}")[-1] if isinstance(elem.tag, str) else "node"
                parts.append(f"{tag_name}: {text}")
        return "\n".join(parts).strip() or raw

    def _parse_yaml(self, file_path: str) -> str:
        raw = self._parse_text(file_path)
        if yaml is None:
            return raw
        try:
            data = yaml.safe_load(raw)
        except Exception:
            return raw
        return json.dumps(data, ensure_ascii=False, indent=2)

    def _parse_html(self, file_path: str) -> str:
        raw = self._parse_text(file_path)
        soup = BeautifulSoup(raw, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        text = soup.get_text("\n")
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        return "\n".join(lines)

    def _parse_rtf(self, file_path: str) -> str:
        raw = self._parse_text(file_path)
        if rtf_to_text is None:
            return raw
        try:
            return rtf_to_text(raw)
        except Exception:
            return raw

    def _parse_image(self, file_path: str) -> str:
        if RapidOCR is None:
            return ""

        try:
            ocr_engine = RapidOCR()
            result = ocr_engine(file_path)
        except Exception:
            return ""

        if isinstance(result, tuple):
            result = result[0]

        if not result:
            return ""

        lines: list[str] = []
        for item in result:
            text = ""
            if isinstance(item, (list, tuple)):
                if len(item) >= 2:
                    text = str(item[1]).strip()
                elif len(item) == 1:
                    text = str(item[0]).strip()
            elif isinstance(item, dict):
                text = str(
                    item.get("text")
                    or item.get("transcription")
                    or item.get("value")
                    or ""
                ).strip()
            else:
                text = str(item).strip()

            if text:
                lines.append(text)

        return "\n".join(lines).strip()

    def _normalize_cell_value(self, value) -> str:
        if value is None:
            return ""
        if hasattr(value, "text"):
            return str(getattr(value, "text") or "").strip()
        if isinstance(value, float) and value.is_integer():
            return str(int(value))
        return str(value).strip()

    def _format_table_rows(self, rows: list[list[str]]) -> str:
        formatted_rows = []
        for row in rows:
            cleaned = [self._normalize_cell_value(cell) for cell in row]
            if any(cleaned):
                formatted_rows.append("\t".join(cleaned).rstrip())
        return "\n".join(formatted_rows).strip()

    def _parse_text(self, file_path: str) -> str:
        encodings = ["utf-8-sig", "utf-8", "gb18030", "utf-16", "latin-1"]
        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding, errors="strict") as f:
                    return f.read()
            except Exception:
                continue
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    def chunk_text(
        self, text: str, chunk_size: int = 1000, overlap: int = 200
    ) -> List[str]:
        """Split text into chunks with overlap"""
        chunks = []
        start = 0
        text_len = len(text)

        while start < text_len:
            end = min(start + chunk_size, text_len)
            chunk = text[start:end]
            chunks.append(chunk)
            start += chunk_size - overlap

        return chunks


document_service = DocumentService()
